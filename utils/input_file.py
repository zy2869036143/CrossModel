import torch
from PIL import Image

import cn_clip.clip as clip
from cn_clip.clip import load_from_name, available_models
from io import BytesIO
import sys
import importlib

importlib.reload(sys)
import codecs

sys.stdout = codecs.getwriter("utf-8")(sys.stdout.detach())
import os
import pandas as pd
import pdfminer
from pdfminer.pdfparser import PDFParser
from pdfminer.pdfdocument import PDFDocument
from pdfminer.pdfpage import PDFPage, PDFTextExtractionNotAllowed
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.layout import LAParams
from pdfminer.converter import PDFPageAggregator
from pptx import Presentation
import aspose.slides as slides
from docx import Document
import jieba
from collections import Counter
from gensim import corpora
import re


# 读取文档路径下的文档名列表
def read_name(_path):
    """读取文档路径下的文档名列表
    :param _path: 文档路径
    :return: 文档名列表
    """
    name_list_temp = []
    files = os.listdir(_path)  # 得到文件夹下的所有一级子路径名称
    for file in files:  # 遍历路径筛选文件
        if '.' in file:
            name_list_temp.append(file)
    return name_list_temp


def remove_en_blank(sentence):
    """
    清除语句块中的英文和空格
    :param sentence: 语句块
    :return: 清除了英文和空格的语句块
    """
    relu = re.compile(r'[ a-zA-Z]')
    relu1 = re.compile(r'[0-9]*')
    relu2 = re.compile(r'[^A-Z^a-z^0-9^\u4e00-\u9fa5]')
    relu3 = re.compile(r'\\t|\\n')
    res = relu.sub('', sentence)
    res = relu1.sub('', str(res))
    res = relu2.sub('', res)
    res = relu3.sub('', str(res))
    return str(res)


# 读取pptx文件
def read_from_pptx(path):
    """
    读取pptx文件
    :param path: 文件路径
    :return: 文件内容
    """
    ppt = Presentation(path)
    context = ''
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text:
                        context += remove_en_blank(paragraph.text.strip())
    return context

def read_from_pptx_text_image(path):
    results = []
    context = ''
    images = []
    ppt = Presentation(path)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text:
                        context += remove_en_blank(paragraph.text.strip())
            # elif isinstance(shape, Picture):
            else:
                try:
                    imdata = shape.image.blob
                    images.append(Image.open(BytesIO(imdata)))
                except:
                    pass
    return context, images


def read_from_pdf(path):
    """
    读取PDF文件
    :param path: 文件路径
    :return: 文件内容
    """
    context = ''
    fp = open(path, 'rb')
    parser = PDFParser(fp)
    document = PDFDocument(parser)
    if not document.is_extractable:
        raise PDFTextExtractionNotAllowed
    rsrcmgr = PDFResourceManager()
    laparams = LAParams(
        char_margin=10.0,
        line_margin=0.2,
        boxes_flow=0.2,
        all_texts=False,
    )
    device = PDFPageAggregator(rsrcmgr, laparams=laparams)
    interpreter = PDFPageInterpreter(rsrcmgr, device)
    for page in PDFPage.create_pages(document):
        interpreter.process_page(page)
        layout = device.get_result()
        for obj in layout._objs:
            if isinstance(obj, pdfminer.layout.LTTextBoxHorizontal):
                context += obj.get_text()
    return context


def read_from_xlsx(path):
    """
    读取xlsx文件
    :param path: 文件路径
    :return: 文件内容
    """
    df = pd.read_excel(path)
    data = df.values
    return data


def read_from_ppt(path):
    """
    读取ppt文件
    :param path: 文件路径
    :return: 文件内容
    """
    text = ''
    images = []
    try:
        pres = slides.Presentation(path)
        newpath = os.path.splitext(path)[0] + '.pptx'
        pres.save(newpath, slides.export.SaveFormat.PPTX)
        text, images = read_from_pptx(newpath)
        os.remove(newpath)
    except Exception as e:
        print(e)
    return text, images


def read_from_docx(path):
    """
    读取docx文件
    :param path: 文件路径
    :return: 文件内容
    """
    content = ''
    _doc = Document(path)
    for _p in _doc.paragraphs:
        content += remove_en_blank(_p.text.strip())
    return content


def read_from_txt(pa):
    """
    读取txt文件
    :param path: 文件路径
    :return: 文件内容
    """
    f = open(pa, "r", encoding="utf-8")
    content = f.read()
    content = remove_en_blank(content)
    f.close()
    return content

def read_from_photo(pa):
    """
    读取照片
    :param path: 文件路径
    :return: 文件内容
    """
    return [Image.open(pa)]


def read(path, file):
    """
    读取文档内容
    :param path: 文件路径
    :param file: 文件名
    :return: 文件内容
    """
    text = ''
    images = []
    file_type = file.split('.')[1]
    print(path + file)
    try:
        if file_type == 'pptx':
            text, images = read_from_pptx_text_image(path + file)
        elif file_type == 'pdf':
            text = read_from_pdf(path + file)
        elif file_type == 'xlsx':
            text = read_from_xlsx(path + file)
        elif file_type == 'ppt':
            text, images = read_from_ppt(path + file)
        elif file_type == 'docx':
            text = read_from_docx(path + file)
        elif file_type == 'txt':
            text = read_from_txt(path + file)
        # elif file_type == 'png' or 'jpeg' or 'jpg':
        #     images = read_from_photo(path + file)
    except:
        return ''
    else:
        text = remove_en_blank(str(text))
        return [text], images


def get_image_features(model, preprocess, pil_images, device):
    if len(pil_images) == 0:
        return None
    print()
    if type(pil_images).__name__ == "list":
        image_feature = model.encode_image(preprocess(pil_images[0]).unsqueeze(0).to(device))
        image_feature /= image_feature.norm(dim=-1, keepdim=True)
        image_features = image_feature
        for pil_image in pil_images[1:]:
            image_feature = model.encode_image(preprocess(pil_image).unsqueeze(0).to(device))
            image_feature /= image_feature.norm(dim=-1, keepdim=True)
            image_features = torch.cat([image_features, image_feature])
    else:
        image_feature = model.encode_image(preprocess(pil_images).unsqueeze(0).to(device))
        image_feature /= image_feature.norm(dim=-1, keepdim=True)
        image_features = image_feature

    return image_features


def get_text_feature(model, text, device):
    token = clip.tokenize(text).to(device)
    text_features, text_features_gloabl = model.encode_text(token)
    print(text_features.shape, text_features_gloabl.shape)
    # 对特征进行归一化，请使用归一化后的图文特征用于下游任务
    text_features_gloabl /= text_features_gloabl.norm(dim=-1, keepdim=True)
    return text_features_gloabl


def get_feature_logits(model, feature1, feature2):
    # cosine similarity as logits
    logit_scale = model.logit_scale.exp()
    logits = logit_scale * feature1 @ feature2.t()
    return logits


def probe(logit):
    return logit.softmax(dim=-1).cpu().numpy()


def get_model():
    print("Available models:", available_models())
    # Available models: ['ViT-B-16', 'ViT-L-14', 'ViT-L-14-336', 'ViT-H-14', 'RN50']

    device = "cuda" if torch.cuda.is_available() else "cpu"
    model, preprocess = load_from_name("ViT-B-16", device=device, download_root='./')
    model.eval()
    return model, preprocess, device

def get_cosine(features_predict, features_tag):
    sim = []
    for feature_predict in features_predict:
        temp = []
        for feature in features_tag:
            similarity = torch.cosine_similarity(feature_predict, feature, dim=0)
            temp.append(similarity)
        sim.append(temp)

    return sim

if __name__ == '__main__':
    print(read_name("F:\lab\项目实训-多级分类标签\数据\文化（包括网络文学、文学名著）\\"))
    doc_path = "F:\lab\项目实训-多级分类标签\数据\文化（包括网络文学、文学名著）\\"
    doc_name = "哈士奇.jpeg"
    #doc_name = doc_name.decode('utf-8')
    #names = read_name(doc_path)

    all_doc_text, images = read(doc_path,doc_name)
    print(all_doc_text, images)

    model, preprocess, device = get_model()
    image_features = get_image_features(model, preprocess, images, device)
    print(image_features.shape)
    tag_features = get_text_feature(model, ["杰尼龟", "妙蛙种子", "小火龙", "皮卡丘", "狗"], device)
    print(get_cosine(image_features, tag_features))
    text_test_features = get_text_feature(model, all_doc_text, device)
    print("tag_features", tag_features.shape)
    print("text_test_features.shape",text_test_features.shape)
    with torch.no_grad():
        logit = get_feature_logits(model, image_features, tag_features)
        probs = probe(logit)

    print("Label probs:", probs)  # [[1.268734e-03 5.436878e-02 6.795761e-04 9.436829e-01]]
